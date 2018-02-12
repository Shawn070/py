from pptx import Presentation

prs = Presentation()    # 创建一个PPT文档
title_slide_layout = prs.slide_layouts[0]   # .slide_layouts[]：确定幻灯片的先后顺序
slide = prs.slides.add_slide(title_slide_layout)    #.slides.add_slide：增加一个幻灯片
title = slide.shapes.title      # 表示一个幻灯片的标题
subtitle = slide.placeholders[1]    # 表示一个幻灯片的内容

title.text = "Hello, World!"
subtitle.text = "python-pptx was here!"

prs.save("pptxTry.pptx")