from pptx import Presentation

prs = Presentation()
bullet_slide_layout = prs.slide_layouts[1]  # 确定顺序

slide = prs.slides.add_slide(bullet_slide_layout)   # 增加幻灯片
shapes = slide.shapes   # shapes：对幻灯片的某一区域操作

title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = "Adding a Bullet Slide"
tf = body_shape.text_frame
tf.text = "Find the bullet slide layout"

p = tf.add_paragraph()  # 增加文本
p.text = "Use __TextFrame.text for first bullet"
p.level = 1             # 设置级别

p = tf.add_paragraph()
p.text = "Use __TextFrame.add_paragraph() for subsequent bullets"
p.level = 2

prs.save("pptxTry_text.pptx")

'''
shapes.add_textbox()
shapes.add_picture()
shapes.add_shape()
shapes.add_table()
'''
